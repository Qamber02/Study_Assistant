import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
from pathlib import Path


# Supported plain-text extensions
TEXT_EXTENSIONS = {".txt", ".py", ".java", ".cpp", ".js", ".ts", ".md", ".csv"}


def _extract_pdf(file_path: Path) -> str:
    with fitz.open(file_path) as doc:
        return "".join(page.get_text() for page in doc)


def _extract_docx(file_path: Path) -> str:
    doc = Document(file_path)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())


def _extract_pptx(file_path: Path) -> str:
    prs = Presentation(file_path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def _extract_text_file(file_path: Path) -> str:
    # Try UTF-8 first, fall back to latin-1 to avoid crashes on weird encodings
    for encoding in ("utf-8", "latin-1"):
        try:
            return file_path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    raise ValueError(f"Could not decode file: {file_path}")


# Dispatch table — easy to extend without touching extract_text()
_EXTRACTORS = {
    ".pdf":  _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
}


def extract_text(file_path: str | Path) -> str:
    """
    Extract text from a file.

    Supports: PDF, DOCX, PPTX, and common plain-text formats.

    Args:
        file_path: Path to the file (str or Path object).

    Returns:
        Extracted text as a string.

    Raises:
        FileNotFoundError: If the file doesn't exist.
        ValueError: If the format is unsupported or file can't be decoded.
        RuntimeError: If extraction fails for any other reason.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"File not found: {path}")

    ext = path.suffix.lower()

    try:
        if ext in _EXTRACTORS:
            return _EXTRACTORS[ext](path)
        elif ext in TEXT_EXTENSIONS:
            return _extract_text_file(path)
        else:
            raise ValueError(f"Unsupported file format: '{ext}'")

    except (FileNotFoundError, ValueError):
        raise  # re-raise known errors as-is
    except Exception as e:
        raise RuntimeError(f"Failed to extract text from '{path.name}': {e}") from e
